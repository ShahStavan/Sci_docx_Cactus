import pytesseract
from pdf2image import convert_from_path
from hotpdf import HotPdf
import fitz
import tabula
import re
import textwrap
from spellchecker import SpellChecker
from docparser import parse 
from pylatexenc.latex2text import LatexNodes2Text
from pptx import Presentation
pytesseract.pytesseract.tesseract_cmd = r'/usr/bin/tesseract'
poppler_path = r'/usr/bin/'

class TextPreprocessor:
    def _init_(self):
        self.spell = SpellChecker()
        self.contractions = {
            "ain't": "am not / are not / is not / has not / have not",
            "aren't": "are not / am not",
            "can't": "cannot",
            "can't've": "cannot have",
            "'cause": "because",
            "could've": "could have",
            "couldn't": "could not",
            "couldn't've": "could not have",
            "didn't": "did not",
            "doesn't": "does not",
            "don't": "do not",
            "hadn't": "had not",
            "hadn't've": "had not have",
            "hasn't": "has not",
            "haven't": "have not",
            "he'd": "he had / he would",
            "he'd've": "he would have",
            "he'll": "he shall / he will",
            "he'll've": "he shall have / he will have",
            "he's": "he has / he is",
            "how'd": "how did",
            "how'd'y": "how do you",
            "how'll": "how will",
            "how's": "how has / how is / how does",
            "i'd": "I had / I would",
            "i'd've": "I would have",
            "i'll": "I shall / I will",
            "i'll've": "I shall have / I will have",
            "i'm": "I am",
            "i've": "I have",
            "isn't": "is not",
            "it'd": "it had / it would",
            "it'd've": "it would have",
            "it'll": "it shall / it will",
            "it'll've": "it shall have / it will have",
            "it's": "it has / it is",
            "let's": "let us",
            "ma'am": "madam",
            "mayn't": "may not",
            "might've": "might have",
            "mightn't": "might not",
            "mightn't've": "might not have",
            "must've": "must have",
            "mustn't": "must not",
            "mustn't've": "must not have",
            "needn't": "need not",
            "needn't've": "need not have",
            "o'clock": "of the clock",
            "oughtn't": "ought not",
            "oughtn't've": "ought not have",
            "shan't": "shall not",
            "sha'n't": "shall not",
            "shan't've": "shall not have",
            "she'd": "she had / she would",
            "she'd've": "she would have",
            "she'll": "she shall / she will",
            "she'll've": "she shall have / she will have",
            "she's": "she has / she is",
            "should've": "should have",
            "shouldn't": "should not",
            "shouldn't've": "should not have",
            "so've": "so have",
            "so's": "so as / so is",
            "that'd": "that would / that had",
            "that'd've": "that would have",
            "that's": "that has / that is",
            "there'd": "there had / there would",
            "there'd've": "there would have",
            "there's": "there has / there is",
            "they'd": "they had / they would",
            "they'd've": "they would have",
            "they'll": "they shall / they will",
            "they'll've": "they shall have / they will have",
            "they're": "they are",
            "they've": "they have",
            "to've": "to have",
            "wasn't": "was not",
            "we'd": "we had / we would",
            "we'd've": "we would have",
            "we'll": "we will",
            "we'll've": "we will have",
            "we're": "we are",
            "we've": "we have",
            "weren't": "were not",
            "what'll": "what shall / what will",
            "what'll've": "what shall have / what will have",
            "what're": "what are",
            "what's": "what has / what is",
            "what've": "what have",
            "when's": "when has / when is",
            "when've": "when have",
            "where'd": "where did",
            "where's": "where has / where is",
            "where've": "where have",
            "who'll": "who shall / who will",
            "who'll've": "who shall have / who will have",
            "who's": "who has / who is",
            "who've": "who have",
            "why's": "why has / why is",
            "why've": "why have",
            "will've": "will have",
            "won't": "will not",
            "won't've": "will not have",
            "would've": "would have",
            "wouldn't": "would not",
            "wouldn't've": "would not have",
            "y'all": "you all",
            "y'all'd": "you all would",
            "y'all'd've": "you all would have",
            "y'all're": "you all are",
            "y'all've": "you all have",
            "you'd": "you had / you would",
            "you'd've": "you would have",
            "you'll": "you shall / you will",
            "you'll've": "you shall have / you will have",
            "you're": "you are",
            "you've": "you have"
            # Add more contractions here...
        }

    def correct_spelling(self, text):
        corrected_text = []
        misspelled_words = self.spell.unknown(text.split())
        for word in text.split():
            next_word = word
            if word in misspelled_words:
                candidates = self.spell.candidates(word)
                next_word = candidates[0] if candidates else word
            corrected_text.append(next_word)
        corrected_text = " ".join(corrected_text)
        wrapped_text = textwrap.fill(corrected_text, width=width)
        return wrapped_text

    def cont_to_exp(self, text):
      if type(text) is str:
          for key in self.contractions:
              value = contractions[key]
              text = text.replace(key, value)
          return text
      else:
          return text

    def remove_urls_and_emails(self, text):
        text_without_urls = re.sub(r'http\S+|www\S+|\S+\.com\S+|\S+\.in\S+', '', text)
        text_without_emails = re.sub(r'\S+@\S+', '', text_without_urls)
        return text_without_emails

    def remove_emojis(self,text):
        emojis_pattern = re.compile(pattern="["
                        u"\U0001F600-\U0001F64F"  # emoticons
                        u"\U0001F300-\U0001F5FF"  # symbols & pictographs
                        u"\U0001F680-\U0001F6FF"  # transport & map symbols
                        u"\U0001F1E0-\U0001F1FF"  # flags (iOS)
                        u"\U00002500-\U00002BEF"  # chinese char
                        u"\U00002702-\U000027B0"
                        u"\U00002702-\U000027B0"
                        u"\U000024C2-\U0001F251"
                        u"\U0001f926-\U0001f937"
                        u"\U00010000-\U0010ffff"
                        u"\u2640-\u2642"
                        u"\u2600-\u2B55"
                        u"\u200d"
                        u"\u23cf"
                        u"\u23e9"
                        u"\u231a"
                        u"\ufe0f"  # dingbats
                        u"\u3030"
                    "]+", flags = re.UNICODE)
        return emojis_pattern.sub(r'', text)

class PdfExtractor:

    def pdf_to_text(self, pdf_path):
        pytesseract.pytesseract.tesseract_cmd = pytesseract_path
        images = convert_from_path(pdf_path, 500, poppler_path=poppler_path)
        text = ''
        for image in images:
            text += pytesseract.image_to_string(image)
        return text

    def hotPdf_to_Text(self, pdf_path):
        hotpdf_document = HotPdf(pdf_path)
        num_pages = len(hotpdf_document.pages)
        full_text = ""
        for page_num in range(num_pages):
            full_page_text = hotpdf_document.extract_page_text(page=page_num)
            full_text += full_page_text + "\n"
        return full_text

    def extract_tables_from_pdf(self, pdf_path):
        tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
        return tables

    def extract_text_from_all_pages(self, pdf_path):
        pdf_document = fitz.open(pdf_path)
        total_pages = pdf_document.page_count
        full_text_by_page = []
        for page_number in range(total_pages):
            page = pdf_document[page_number]
            page_text = page.get_text("text")
            full_text_by_page.append(page_text)
        pdf_document.close()
        return full_text_by_page

    def extract_tables_and_text(self, pdf_path):
        tables = self.extract_tables_from_pdf(pdf_path)
        text = self.extract_text_from_all_pages(pdf_path)
        combined_data = {"tables": tables, "text": text}
        return combined_data

class DocxExtractor:
    def _init_(self):
        pass

    def docx_to_text(self, docx_path):
        # parse = pydocparser.Parser()
      document = parse(doc_path)
      content = document.content

      # Split the text into lines, remove leading and trailing spaces, and join again
      formatted_text = ' '.join(line.strip() for line in content.splitlines())

      # Replace multiple consecutive spaces with a single space
      formatted_text = re.sub(r'\s+', ' ', formatted_text)

      # Set the desired width for the structured print (large value to output the whole text)
      width = 80

      # Wrap the text to the specified width (which will output the whole text)
      wrapped_text = textwrap.fill(formatted_text, width=width)

      # Print the structured text
      return wrapped_text

class TexExtractor:
    def _init_(self):
        pass

    def latex_file_to_plaintext(self, latex_file_path):
        try:
          with open(latex_file_path, 'r', encoding='utf-8') as f:
              latex_string = f.read()
          latex_string = re.sub(r'\\begin\{tabular\}\{(.?)\}(.?)\\end\{tabular\}', lambda match: parse_tabular_latex(match), latex_string, flags=re.DOTALL)
          latex_string = re.sub(r'\\includegraphics(?:\[[^\]]\])?\{[^\}]\}', '', latex_string)
          text = LatexNodes2Text().latex_to_text(latex_string)
          # text = re.sub(r'[^\w\s]', '', text)
          text = text.lower()
          # text = re.sub(r'[^a-zA-Z0-9\s+\-*=/\\\[\]{}()^_]', '', text)
          text = re.sub(r'\n{4,}', '\n\n\n', text)
          text = text.replace("§","")
          return text
        except Exception as e:
            print("Conversion failed:", e)
            return None
class PptxExtractor:
    def _init_(self):
        pass

    def extract_text_from_ppt(ppt_path):
        prs = Presentation(ppt_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
        return text

# pdf_extractor = PdfExtractor()
# pdf_text = pdf_extractor.pdf_to_text("example.pdf")
# print(pdf_text)

# docx_extractor = DocxExtractor()
# docx_text = docx_extractor.docx_to_text("example.docx")
# print(docx_text)